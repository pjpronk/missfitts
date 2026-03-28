from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from lxml import etree
from PIL import Image, ImageDraw
import io

# --- Colors ---
BG        = RGBColor(0x0b, 0x15, 0x0a)
GREEN_DIM = RGBColor(0x3B, 0x6D, 0x11)
GREEN_MID = RGBColor(0x63, 0x99, 0x22)
GREEN_LT  = RGBColor(0xC0, 0xDD, 0x97)
GREEN_BAR = RGBColor(0x1a, 0x2d, 0x14)
GREEN_BRT = RGBColor(0x97, 0xC4, 0x59)  # bright mid — bullets + subtitles

W = Inches(13.33)
H = Inches(7.5)

A_NS = 'http://schemas.openxmlformats.org/drawingml/2006/main'
P_NS = 'http://schemas.openxmlformats.org/presentationml/2006/main'

# ── Helpers ───────────────────────────────────────────────────────────────────

def lock_shape(shape):
    sp = shape._element
    nvSpPr = sp.find(f'.//{{{P_NS}}}nvSpPr')
    if nvSpPr is None:
        nvSpPr = sp.find('.//{http://schemas.openxmlformats.org/drawingml/2006/spreadsheetDrawing}nvSpPr')
    cNvSpPr = sp.find(f'.//{{{P_NS}}}cNvSpPr')
    if cNvSpPr is not None:
        locks = etree.SubElement(cNvSpPr, f'{{{A_NS}}}spLocks')
        locks.set('noSelect', '1')
        locks.set('noResize', '1')
        locks.set('noMove', '1')

def add_rect(slide, x, y, w, h, fill=None, line_color=None, line_width=Pt(0), locked=False):
    shape = slide.shapes.add_shape(1, x, y, w, h)
    shape.line.width = line_width
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
    else:
        shape.line.fill.background()
    if locked:
        lock_shape(shape)
    return shape

def add_text(slide, text, x, y, w, h, size, color, bold=False,
             align=PP_ALIGN.LEFT, font_name="Courier New", locked=False):
    txb = slide.shapes.add_textbox(x, y, w, h)
    tf = txb.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = size
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.name = font_name
    if locked:
        lock_shape(txb)
    return txb

def placeholder_image(w_inches, h_inches, label="IMAGE"):
    """Generate a dark green placeholder image as a BytesIO PNG."""
    px_w = int(w_inches * 96)
    px_h = int(h_inches * 96)
    img = Image.new("RGB", (px_w, px_h), (17, 34, 11))
    draw = ImageDraw.Draw(img)
    # Border
    draw.rectangle([2, 2, px_w-3, px_h-3], outline=(99, 153, 34), width=2)
    # Diagonal crosses
    draw.line([0, 0, px_w, px_h], fill=(40, 80, 20), width=1)
    draw.line([px_w, 0, 0, px_h], fill=(40, 80, 20), width=1)
    # Label
    draw.text((px_w//2 - 30, px_h//2 - 8), label, fill=(99, 153, 34))
    buf = io.BytesIO()
    img.save(buf, format="PNG")
    buf.seek(0)
    return buf

# ── Corner brackets ───────────────────────────────────────────────────────────

CORNER_T = Inches(0.055)   # thickness — adjust this to taste
CORNER_S = Inches(0.35)    # arm length
PAD      = Inches(0.35)    # distance from slide edge

def draw_corners(slide):
    """Draw all 4 corner brackets, locked."""
    t, s, p = CORNER_T, CORNER_S, PAD

    # Top-left
    add_rect(slide, p, p, s, t, fill=GREEN_MID, locked=True)
    add_rect(slide, p, p, t, s, fill=GREEN_MID, locked=True)
    # Top-right
    add_rect(slide, W - p - s, p, s, t, fill=GREEN_MID, locked=True)
    add_rect(slide, W - p - t, p, t, s, fill=GREEN_MID, locked=True)
    # Bottom-left
    add_rect(slide, p, H - p - t, s, t, fill=GREEN_MID, locked=True)
    add_rect(slide, p, H - p - s, t, s, fill=GREEN_MID, locked=True)
    # Bottom-right
    add_rect(slide, W - p - s, H - p - t, s, t, fill=GREEN_MID, locked=True)
    add_rect(slide, W - p - t, H - p - s, t, s, fill=GREEN_MID, locked=True)

# ── Shared chrome (locked) ────────────────────────────────────────────────────

def draw_hud_chrome(slide, idx, total, progress_label=None):
    pad = PAD
    add_rect(slide, 0, 0, W, H, fill=BG, locked=True)
    draw_corners(slide)
    # Accent line + title divider
    add_rect(slide, pad + Inches(0.15), Inches(1.2), Inches(0.5), Inches(0.05), fill=GREEN_MID, locked=True)
    add_rect(slide, pad + Inches(0.15), Inches(2.05),
             W - pad * 2 - Inches(0.3), Inches(0.02), fill=GREEN_DIM, locked=True)
    # Progress bar — pulled in from corners, taller, more prominent
    bar_x = pad + Inches(0.6)
    bar_w = W - pad * 2 - Inches(1.2)
    bar_y = H - pad - Inches(0.28)
    bar_h = Inches(0.08)
    add_rect(slide, bar_x, bar_y, bar_w, bar_h, fill=GREEN_BAR, locked=True)
    pct = 0 if idx < 0 else round((idx + 1) / total * 100)
    fill_w = 0 if idx < 0 else int(bar_w * (idx + 1) / total)
    add_rect(slide, bar_x, bar_y, fill_w, bar_h, fill=GREEN_MID)
    label = progress_label if progress_label else f"PROGRESS: {pct}%"
    add_text(slide, label, bar_x, bar_y - Inches(0.28),
             Inches(3), Inches(0.26), Pt(9), GREEN_BRT)

def draw_header(slide, tag, num, title, subtitle):
    pad = PAD
    add_text(slide, tag.upper(), pad + Inches(0.3), pad - Inches(0.02),
             Inches(5), Inches(0.3), Pt(10), GREEN_DIM)
    add_text(slide, num, W - pad - Inches(1.6), pad - Inches(0.02),
             Inches(1.6), Inches(0.3), Pt(10), GREEN_DIM, align=PP_ALIGN.RIGHT)
    add_text(slide, title.upper(), pad + Inches(0.15), Inches(1.3),
             W - pad * 2 - Inches(0.3), Inches(0.7), Pt(28), GREEN_LT, bold=True)
    add_text(slide, subtitle.upper(), pad + Inches(0.15), Inches(2.1),
             W - pad * 2, Inches(0.4), Pt(11), GREEN_MID)

# ── Slide types ───────────────────────────────────────────────────────────────

def make_toc_slide(prs, items):
    """items = list of (num, label) e.g. [('01', 'Introduction'), ...]"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    pad = PAD
    bar_x = pad + Inches(0.6)
    bar_w = W - pad * 2 - Inches(1.2)
    bar_y = H - pad - Inches(0.28)
    bar_h = Inches(0.08)

    draw_hud_chrome(slide, -1, total)
    draw_header(slide, "Index", "00 / 07", "Table of Contents", "")

    # Boot header
    add_text(slide, ">  LOADING MISSION INDEX........... OK", bar_x, Inches(2.15),
             bar_w, Inches(0.3), Pt(8), GREEN_DIM)

    # TOC items
    col_w = bar_w / 2 - Inches(0.2)
    for i, (num, label) in enumerate(items):
        col = i % 2
        row = i // 2
        x = bar_x + col * (bar_w / 2)
        y = Inches(2.6) + row * Inches(0.55)
        add_text(slide, f">  {num}  //  {label.upper()}", x, y,
                 col_w, Inches(0.4), Pt(11), GREEN_BRT)

    # Ready line
    add_text(slide, ">  AWAITING INPUT _", bar_x, Inches(5.9),
             bar_w, Inches(0.3), Pt(8), GREEN_DIM)

    # Progress bar at 0%
    add_rect(slide, bar_x, bar_y, bar_w, bar_h, fill=GREEN_BAR, locked=True)
    add_text(slide, "PROGRESS: 0%", bar_x, bar_y - Inches(0.28),
             Inches(3), Inches(0.26), Pt(9), GREEN_BRT)


def make_bullet_slide(prs, tag, num, title, subtitle, bullets, idx, total, progress_label=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    draw_hud_chrome(slide, idx, total, progress_label)
    draw_header(slide, tag, num, title, subtitle)
    pad = PAD
    for i, bullet in enumerate(bullets):
        add_text(slide, f"> {bullet}", pad + Inches(0.15),
                 Inches(2.65) + Inches(0.48) * i,
                 W - pad * 2 - Inches(0.3), Inches(0.4), Pt(13), GREEN_MID)

def make_one_image_slide(prs, tag, num, title, subtitle, idx, total, progress_label=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    draw_hud_chrome(slide, idx, total, progress_label)
    draw_header(slide, tag, num, title, subtitle)
    pad = PAD
    img_x = pad + Inches(0.15)
    img_y = Inches(2.4)
    img_w = W - pad * 2 - Inches(0.3)
    img_h = Inches(3.6)
    buf = placeholder_image(img_w / 914400, img_h / 914400, "[ INSERT IMAGE ]")
    slide.shapes.add_picture(buf, img_x, img_y, img_w, img_h)

def make_two_image_slide(prs, tag, num, title, subtitle, idx, total, progress_label=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    draw_hud_chrome(slide, idx, total, progress_label)
    draw_header(slide, tag, num, title, subtitle)
    pad = PAD
    gap = Inches(0.2)
    img_y = Inches(2.4)
    img_h = Inches(3.5)
    img_w = (W - pad * 2 - Inches(0.3) - gap) / 2
    for i in range(2):
        img_x = pad + Inches(0.15) + i * (img_w + gap)
        label = "[ IMAGE 1 ]" if i == 0 else "[ IMAGE 2 ]"
        buf = placeholder_image(img_w / 914400, img_h / 914400, label)
        slide.shapes.add_picture(buf, img_x, img_y, img_w, img_h)

def make_title_slide(prs, title, subtitle, image_path=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    pad = PAD

    # Background
    add_rect(slide, 0, 0, W, H, fill=BG, locked=True)
    draw_corners(slide)

    # Progress bar track (0% on title slide)
    bar_y = H - pad - Inches(0.18)
    bar_x = pad + Inches(0.3)
    bar_w = W - pad * 2 - Inches(0.6)
    add_rect(slide, bar_x, bar_y, bar_w, Inches(0.055), fill=GREEN_BAR, locked=True)
    add_text(slide, "PROGRESS: 0%", bar_x, bar_y - Inches(0.24),
             Inches(2), Inches(0.25), Pt(8), GREEN_DIM, locked=True)

    # Top labels
    add_text(slide, "TITLE SLIDE", pad + Inches(0.3), pad - Inches(0.02),
             Inches(4), Inches(0.3), Pt(10), GREEN_DIM, locked=True)
    add_text(slide, "00 / 07", W - pad - Inches(1.6), pad - Inches(0.02),
             Inches(1.6), Inches(0.3), Pt(10), GREEN_DIM, align=PP_ALIGN.RIGHT, locked=True)

    # Vertical divider at 48% width
    div_x = W * 0.48
    add_rect(slide, div_x, 0, Inches(0.02), H, fill=GREEN_DIM, locked=True)

    # Right image — fills full right half edge to edge
    img_x = div_x + Inches(0.02)
    img_w = W - img_x
    if image_path:
        slide.shapes.add_picture(image_path, img_x, Emu(0), img_w, H)
    else:
        add_rect(slide, img_x, Emu(0), img_w, H,
                 fill=RGBColor(0x06, 0x0e, 0x04),
                 line_color=GREEN_DIM, line_width=Pt(1))

    tx = pad + Inches(0.3)
    left_w = div_x - pad - Inches(0.3)

    # Boot lines — same key/value style as meta
    boot = [
        (">  SYSTEM BOOT..........", "OK"),
        (">  LOADING ASSETS.......", "OK"),
        (">  ESTABLISHING LINK....", "OK"),
    ]
    for i, (key, val) in enumerate(boot):
        add_text(slide, f"{key} {val}", tx, Inches(1.5) + Inches(0.36) * i,
                 left_w, Inches(0.3), Pt(8), GREEN_DIM)

    # Accent rule above title
    add_rect(slide, tx, Inches(2.72), Inches(0.28), Inches(0.04), fill=GREEN_MID)

    # Title
    add_text(slide, title.upper(), tx, Inches(2.82),
             left_w, Inches(0.85), Pt(22), GREEN_LT, bold=True)

    # Full-width rule between title and subtitle
    add_rect(slide, tx, Inches(3.75), left_w, Inches(0.015), fill=GREEN_DIM)

    # Subtitle (bright, clearly readable)
    add_text(slide, subtitle.upper(), tx, Inches(3.82),
             left_w, Inches(0.38), Pt(10), GREEN_LT)

    # Terminal author/meta list
    meta = [
        (">  CALLSIGN..............", "MISSFITS"),
        (">  UNIT...................", "MISSES & FITS LAW"),
        (">  DATE OF BRIEFING......", "26-03-2026"),
        (">  LEGAL AMMO STATUS.....", "FULLY LOADED"),
    ]
    for i, (key, val) in enumerate(meta):
        line = f"{key} {val}"
        add_text(slide, line, tx, Inches(4.35) + Inches(0.36) * i,
                 left_w, Inches(0.3), Pt(8), GREEN_MID)

    add_text(slide, ">  READY _", tx, Inches(5.82),
             left_w, Inches(0.3), Pt(8), GREEN_DIM)


def make_table_slide(prs, tag, num, title, subtitle, idx, total, progress_label=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    draw_hud_chrome(slide, idx, total, progress_label)
    draw_header(slide, tag, num, title, subtitle)
    pad = PAD

    rows, cols = 4, 4
    tbl_x = pad + Inches(0.15)
    tbl_y = Inches(2.4)
    tbl_w = W - pad * 2 - Inches(0.3)
    tbl_h = Inches(3.4)

    table = slide.shapes.add_table(rows, cols, tbl_x, tbl_y, tbl_w, tbl_h).table

    headers = ["CATEGORY", "VALUE A", "VALUE B", "STATUS"]
    rows_data = [
        ["Item 01", "—", "—", "ACTIVE"],
        ["Item 02", "—", "—", "STANDBY"],
        ["Item 03", "—", "—", "OFFLINE"],
    ]

    for ci, h in enumerate(headers):
        cell = table.cell(0, ci)
        cell.fill.solid()
        cell.fill.fore_color.rgb = GREEN_DIM
        p = cell.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = h
        run.font.size = Pt(11)
        run.font.bold = True
        run.font.color.rgb = GREEN_LT
        run.font.name = "Courier New"

    for ri, row in enumerate(rows_data):
        for ci, val in enumerate(row):
            cell = table.cell(ri + 1, ci)
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(0x0f, 0x20, 0x0a) if ri % 2 == 0 else BG
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            run = p.add_run()
            run.text = val
            run.font.size = Pt(11)
            run.font.color.rgb = GREEN_MID
            run.font.name = "Courier New"

# ── Build ─────────────────────────────────────────────────────────────────────

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

total = 8

make_title_slide(prs,
    title="Operation Missfits",
    subtitle="Misses & Fits Law — Mission Briefing",
    image_path=None)          # ← replace None with r"path\to\your\image.jpg"

make_toc_slide(prs, [
    ("01", "Operation Briefing"),
    ("02", "Primary Objectives"),
    ("03", "Visual Reference"),
    ("04", "Dual Comparison"),
    ("05", "Data Table"),
    ("06", "Field Analysis"),
    ("07", "Mission Debrief"),
])

make_bullet_slide(prs, "Mission Brief", "01 / 07", "Operation Briefing",
    "Strategic Overview — Unit Alpha",
    ["Use this slide for your title and mission statement",
     "Subtitle acts as a secondary context line",
     "Corner brackets define the slide boundary"], 0, total)

make_bullet_slide(prs, "Objectives", "02 / 07", "Primary Objectives",
    "Phase 1 Targets — Confirmed",
    ["State your main argument or goal clearly",
     "Each bullet maps to one clear action or point",
     "Color draws the eye to what matters most"], 1, total)

make_one_image_slide(prs, "Visual", "03 / 07", "Single Image Layout",
    "Replace placeholder with your image", 2, total)

make_two_image_slide(prs, "Comparison", "04 / 07", "Dual Image Layout",
    "Side-by-side visual comparison", 3, total)

make_table_slide(prs, "Data", "05 / 07", "Data Table Layout",
    "Edit rows and columns as needed", 4, total)

make_bullet_slide(prs, "Analysis", "06 / 07", "Field Analysis",
    "Data Review — Sector 4",
    ["Add your analysis points here",
     "Keep bullets short and impactful",
     "Replace with your own content"], 5, total)

make_bullet_slide(prs, "Debrief", "07 / 07", "Mission Debrief",
    "Post-Action Summary — All Units",
    ["Summarize key takeaways here",
     "Call to action goes on the last slide",
     "Progress bar reaches 100% — mission complete"], 6, total)

from datetime import datetime
filename = f"hud_template_{datetime.now().strftime('%Y%m%d_%H%M%S')}.pptx"
prs.save(filename)
print(f"Saved: {filename}")